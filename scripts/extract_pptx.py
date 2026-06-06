#!/usr/bin/env python
"""Extract slide text + speaker notes from a .pptx file. Writes UTF-8 file."""
import sys
from pptx import Presentation


def extract(in_path: str, out_path: str) -> None:
    prs = Presentation(in_path)
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"\n## Slide {i}")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                runs = [run.text for run in para.runs if run.text]
                if runs:
                    lines.append("".join(runs))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                lines.append("\n[Notes]")
                lines.append(note)
    with open(out_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))


if __name__ == "__main__":
    extract(sys.argv[1], sys.argv[2])
    print(f"extracted -> {sys.argv[2]}")
