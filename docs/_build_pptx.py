# -*- coding: utf-8 -*-
"""kiduri 紹介 PowerPoint。きづりサイト寄りの温かいデザイン+現場価値ファースト+サイト画像活用。"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__)); FIG = os.path.join(HERE, "figures")
FONT = "BIZ UDPGothic"

CREAM=RGBColor(0xF7,0xF5,0xF0); INK=RGBColor(0x2C,0x3E,0x50); ORANGE=RGBColor(0xE8,0x83,0x3A)
TEAL=RGBColor(0x2A,0x8F,0x8F); NAVY=RGBColor(0x1A,0x27,0x44); BEIGE=RGBColor(0xEF,0xEA,0xE0)
GRAY=RGBColor(0x6B,0x7B,0x8D); WHITE=RGBColor(0xFF,0xFF,0xFF)
CDW=RGBColor(0xC9,0xD6,0xE6); CBW=RGBColor(0xB8,0xC6,0xD8)

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]; TOTAL = 20; SW=13.333; SH=7.5
_PG=[0]  # スライド作成順に自動採番(差し込み時の番号振り直しを不要にする)

def slide(bg=CREAM):
    _PG[0]+=1
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.shadow.inherit=False; r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background()
    return s

def rect(s,x,y,w,h,fill=None,line=None,shape=MSO_SHAPE.RECTANGLE,radius=None):
    sp=s.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h)); sp.shadow.inherit=False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1.5)
    if radius is not None and shape==MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0]=radius
        except Exception: pass
    return sp

def set_alpha(shape,pct):
    srgb=shape._element.spPr.find(qn('a:solidFill')).find(qn('a:srgbClr'))
    al=srgb.makeelement(qn('a:alpha'),{'val':str(int(pct*1000))}); srgb.append(al)

def pic_cover(s,path,x,y,w,h,border=None):
    iw,ih=Image.open(path).size; iar=iw/ih; bar=w/h
    pic=s.shapes.add_picture(path,Inches(x),Inches(y),Inches(w),Inches(h))
    if iar>bar: c=(1-bar/iar)/2; pic.crop_left=c; pic.crop_right=c
    elif iar<bar: c=(1-iar/bar)/2; pic.crop_top=c; pic.crop_bottom=c
    if border: pic.line.color.rgb=border; pic.line.width=Pt(1.25)
    return pic

def _run(p,text,size,color=INK,bold=False,italic=False,name=FONT):
    r=p.add_run(); r.text=text; f=r.font
    f.size=Pt(size); f.bold=bold; f.italic=italic; f.color.rgb=color; f.name=name
    rPr=r._r.get_or_add_rPr(); ea=rPr.makeelement(qn('a:ea'),{'typeface':name}); rPr.append(ea)
    return r

def tb(s,x,y,w,h,anchor=MSO_ANCHOR.TOP):
    box=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=box.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    return tf

def para(tf,first=False,align=PP_ALIGN.LEFT,sa=8,sb=0,level=0,line=1.15):
    p=tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment=align; p.level=level; p.space_after=Pt(sa); p.space_before=Pt(sb); p.line_spacing=line
    return p

def emit(p,text,size,color=INK):
    import re
    for seg in re.split(r'(\*\*[^*]+\*\*)',text):
        if not seg: continue
        if seg.startswith('**'): _run(p,seg[2:-2],size,ORANGE if color==INK else color,bold=True)
        else: _run(p,seg,size,color)

def footer(s,page=None):
    # page引数は後方互換のため残すが、実番号は作成順カウンタ(_PG)を使う
    page=_PG[0]
    rect(s,0,7.16,SW,0.018,fill=RGBColor(0xE3,0xDD,0xD2))
    _run(para(tb(s,0.55,7.2,8,0.28),first=True),"きづり ｜ 株式会社ソーシップ",9,GRAY)
    _run(para(tb(s,11.2,7.2,1.6,0.28),first=True,align=PP_ALIGN.RIGHT),f"{page} / {TOTAL}",9,GRAY)

def title_bar(s,title,sub=None,accent=ORANGE):
    rect(s,0.55,0.5,0.14,0.66,fill=accent,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.5)
    _run(para(tb(s,0.82,0.42,11.9,0.82),first=True,sa=0),title,27,INK,bold=True)
    if sub: _run(para(tb(s,0.84,1.2,11.9,0.4),first=True),sub,13.5,TEAL)
    rect(s,0.82,1.62 if sub else 1.34,11.7,0.016,fill=RGBColor(0xE3,0xDD,0xD2))

def bullets(s,items,x=0.9,y=1.95,w=11.5,h=4.9,size=15.5,gap=10):
    tf=tb(s,x,y,w,h)
    for i,(text,lvl) in enumerate(items):
        p=para(tf,first=(i==0),level=lvl,sa=gap,line=1.2)
        if lvl==0: _run(p,"▸ ",size+1,ORANGE,bold=True); emit(p,text,size)
        else: _run(p,"・",size-1,TEAL); emit(p,text,size-1,RGBColor(0x4A,0x55,0x60))
    return tf

def card(s,x,y,w,h,head,body,accent=ORANGE,fill=WHITE,num=None):
    sp=rect(s,x,y,w,h,fill=fill,line=RGBColor(0xE4,0xDE,0xD3),shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
    tf=sp.text_frame; tf.word_wrap=True
    tf.margin_left=Inches(0.24);tf.margin_right=Inches(0.22);tf.margin_top=Inches(0.22);tf.margin_bottom=Inches(0.18)
    p=tf.paragraphs[0]; p.line_spacing=1.1; p.space_after=Pt(8); p.alignment=PP_ALIGN.LEFT
    if num: _run(p,num+"  ",17,accent,bold=True)
    _run(p,head,15.5,accent,bold=True)
    p2=tf.add_paragraph(); p2.line_spacing=1.26; p2.alignment=PP_ALIGN.LEFT; _run(p2,body,12.5,INK)

def cards_row(s,cards,y=2.1,h=3.5):
    n=len(cards); gap=0.32; w=(12.0-gap*(n-1))/n; x0=0.66
    for i,c in enumerate(cards):
        card(s,x0+i*(w+gap),y,w,h,c[0],c[1],accent=c[2] if len(c)>2 else ORANGE,num=c[3] if len(c)>3 else None)

def image_slide(title,img,caption,page,sub=None):
    s=slide(); title_bar(s,title,sub); footer(s,page)
    path=os.path.join(FIG,img); iw,ih=Image.open(path).size; ar=iw/ih
    bx,by,bw,bh=0.8,1.85,11.73,4.55; w=bw; h=w/ar
    if h>bh: h=bh; w=h*ar
    x=bx+(bw-w)/2; y=by+(bh-h)/2
    rect(s,x-0.05,y-0.05,w+0.1,h+0.1,fill=WHITE,line=RGBColor(0xDD,0xD5,0xC8),shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.02)
    s.shapes.add_picture(path,Inches(x),Inches(y),Inches(w),Inches(h))
    if caption: _run(para(tb(s,0.8,by+bh+0.06,11.73,0.4),first=True,align=PP_ALIGN.CENTER),caption,11,GRAY)
    return s

def content_slide(title,items,page,sub=None,**kw):
    s=slide(); title_bar(s,title,sub); footer(s,page); bullets(s,items,**kw); return s

def cards_slide(title,cards,page,sub=None,y=2.15,h=3.4):
    s=slide(); title_bar(s,title,sub); footer(s,page); cards_row(s,cards,y=y,h=h); return s

def hero_slide(builder):
    s=slide(NAVY); pic_cover(s,os.path.join(FIG,"site_hero.jpg"),0,0,SW,SH)
    ov=rect(s,0,0,SW,SH,fill=NAVY); set_alpha(ov,60)
    builder(s); return s

# ============ S1 タイトル(ヒーロー画像) ============
def s1(s):
    rect(s,0,5.05,SW,0.03,fill=ORANGE)
    _run(para(tb(s,0.95,0.55,9,0.5),first=True),"放課後等デイサービス・児童発達支援事業所向け",13,CDW)
    tf=tb(s,0.9,1.55,11.6,2.7)
    _run(para(tf,first=True,sa=2),"きづり（kiduri）",52,WHITE,bold=True)
    _run(para(tf,sa=0,line=1.2),"書類に追われる日々から、",26,WHITE,bold=True)
    _run(para(tf,sa=0,line=1.2),"子どもと向き合う毎日へ。",26,ORANGE,bold=True)
    tf2=tb(s,0.92,5.3,11.5,1.4)
    _run(para(tf2,first=True,sa=4),"商社・共同開発者向け ご紹介資料",16,WHITE)
    _run(para(tf2,sa=2),"2026年7月1日",12.5,CBW)
    _run(para(tf2,sa=0),"運営：株式会社ソーシップ",12.5,CBW)
hero_slide(s1)

# ============ S2 コンセプト ============
s=slide()
_run(para(tb(s,1.0,1.9,11.3,0.5),first=True),"CONCEPT",14,ORANGE,bold=True)
tf=tb(s,1.0,2.4,11.3,2.0)
_run(para(tf,first=True,sa=6,line=1.2),"「書類のための時間」を、",33,INK,bold=True)
_run(para(tf,sa=0,line=1.2),"「子どものための時間」へ。",33,ORANGE,bold=True)
rect(s,1.05,4.55,3.4,0.06,fill=ORANGE)
_run(para(tb(s,1.05,4.8,11.0,1.5),first=True,line=1.5),
 "毎日の連絡帳の記録が、AIの力で個別支援計画書に変わる。煩雑な事務作業を劇的に軽減し、職員が子どもと向き合う時間を取り戻します。まず現場の困りごとを解決する ―― それが「きづり」の出発点です。",15,INK)
footer(s,2)

# ============ S3 しんどさ ============
cards_slide("その「しんどさ」、放置していませんか？",[
 ("支援後の残業が終わらない","子どもが帰った後、夜遅くまでPCに向かう毎日。書類作成は支援業務の一部なのに、現場の疲弊を加速させています。",ORANGE),
 ("特定のスタッフに業務が集中","計画書を書ける人が限られ負荷が偏る。その人が辞めたら回らない――属人化が経営リスクに直結します。",ORANGE),
 ("期限超過・書類不備の不安","「あの計画書、いつが期限だっけ？」「指導で指摘されたら…」。書類管理の不安が、現場の余裕を奪っています。",ORANGE),
], 3, sub="多くの事業所が抱える共通の課題。ひとつでも心当たりがあれば、変えるタイミングです")

# ============ S4 解決の核 ============
content_slide("「きづり」は、この構造的な課題を根本から変えます", [
 ("毎日の記録が、**アセスメント・モニタリングを通じて計画書づくりの根拠(エビデンス)**になる",0),
 ("計画書の目標は **継続はモニタリングから・新規はアセスメントから**。記録が一貫した根拠として活きる",0),
 ("だから、特別なスキルがなくても **チーム全員で質の高い支援** を実現できる",0),
 ("AIは文章を書く「アシスタント」。**送信・確定は必ず人が確認**(AIが勝手に送らない)",0),
 ("五領域(健康・生活/運動・感覚/認知・行動/言語・コミュニケーション/人間関係・社会性)に対応",0),
], 4, sub="記録 → アセスメント・モニタリング → 計画書を、ひとつの流れにつなぐ")

# ============ S5 機能① ============
cards_slide("現場の1日に、無理なく溶け込む設計",[
 ("朝 ― 支援案を確認","本日の活動目的と出席者を把握。何を見るかを先に共有。",TEAL,"①"),
 ("活動中 ― タブレットで記録","気になった行動・成長をリアルタイムに入力。複数スタッフ同時入力に対応。",ORANGE,"②"),
 ("活動後 ― AIが統合→送信","AIが連絡帳の文章を自動生成。確認・編集してワンクリックで保護者へ(既読確認つき)。",TEAL,"③"),
], 5, sub="ICTが苦手なスタッフでも、スマホ感覚で使える「タブレットモード」を搭載")

# ============ S6 機能②(画像つき) ============
s=slide(); title_bar(s,"6ヶ月周期の書類業務を、AIが自動でナビゲート","「期限、いつだっけ?」をなくす。抜け漏れと不安を仕組みで解消"); footer(s,6)
bullets(s,[
 ("更新期限をシステムが**自動計算**し、**色別アラート**でお知らせ",0),
 ("「きづり(記録シート)」「モニタリング表」「計画書」をAIが**ステップ誘導**",0),
 ("記録がアセスメント・計画へ根拠として引き継がれ、**書き直しの手間が激減**",0),
 ("初期設定・データ移行・操作研修まで専任が伴走",0),
], x=0.9, y=2.0, w=6.7, h=4.6, size=15)
pic_cover(s, os.path.join(FIG,"site_feature_plan.jpg"), 7.95, 2.0, 4.6, 4.0, border=RGBColor(0xDD,0xD5,0xC8))

# ============ S7-8 製品画面 ============
image_slide("製品画面 ― AIが個別支援計画の作成を支援","fig4.png",
 "同意・支援上の特性・問い返しを一画面で。連絡帳の記録から計画書づくりへ(個人情報モザイク済)",7,sub="現場のAI支援")
image_slide("製品画面 ― 施設の記録基準をAIと対話でつくる","fig1.png",
 "管理者がAIと相談しながら施設独自の記録基準を作成。記録の質を施設で揃える(個人情報モザイク済)",8,sub="品質の標準化")

# ============ 計画書の別添「評価状況の全体像」 ============
image_slide("成長を「見える化」する別添 ―「評価状況の全体像」","betten_sample.png",
 "この客観スコアは個別支援計画AIの根拠にも自動で差し込まれ、記録 → 評価 → 計画がひとつながりになる",9,
 sub="支援者の客観 × 本人の主観(mynameis連携)を重ね、計画書の別添PDFとして出力")

# ============ S9 導入事例 ============
s=slide(); title_bar(s,"導入後に生まれる変化","単なる時間消化の支援から、チームで質を追求する文化へ"); footer(s,9)
cases=[("site_case01.jpg","チームで支援を語る時間が生まれる"),
       ("site_case02.jpg","整理された環境で、抜け漏れの不安が減る"),
       ("site_case03.jpg","タブレットで、その場で記録・計画を確認")]
cw=3.74; gap=0.36; x0=0.66; iy=2.0; ih=2.18
for i,(img,cap) in enumerate(cases):
    x=x0+i*(cw+gap)
    pic_cover(s,os.path.join(FIG,img),x,iy,cw,ih,border=RGBColor(0xDD,0xD5,0xC8))
    _run(para(tb(s,x,iy+ih+0.1,cw,0.6),first=True,align=PP_ALIGN.CENTER,line=1.2),cap,11.5,INK)
qb=rect(s,0.9,5.35,11.5,1.25,fill=BEIGE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.08)
tf=qb.text_frame; tf.word_wrap=True
tf.margin_left=Inches(0.3);tf.margin_right=Inches(0.3);tf.margin_top=Inches(0.18)
p=tf.paragraphs[0]; p.line_spacing=1.3
_run(p,"「『この子のために今、何をするか』を意識するように。単なる時間消化の支援から脱却し、チームで支援の質を追求する文化が生まれました。」",13.5,INK)

# ============ S10 導入の流れ ============
cards_slide("導入はかんたん。最短で運用開始まで",[
 ("お問い合わせ・無料相談","施設の状況・お悩みを伺い、最適なプランをご提案。",ORANGE,"01"),
 ("オンラインデモ","実際の画面で、記録から計画書作成までの流れを体験。",TEAL,"02"),
 ("初期設定・データ移行","施設・児童データの登録、既存データ移行を専任が支援。",ORANGE,"03"),
 ("運用開始・定着サポート","操作研修を実施。軌道に乗るまで継続サポート。",TEAL,"04"),
], 10, y=2.3, h=3.0, sub="専任スタッフが初期設定から運用定着まで伴走するので、ICTが苦手な施設でも安心")

# ============ S11 区切り ============
def s11(s):
    rect(s,0.95,2.45,3.0,0.07,fill=ORANGE)
    tf=tb(s,0.95,2.7,11.4,2.2)
    _run(para(tf,first=True,sa=6,line=1.18),"その先の将来設計",30,WHITE,bold=True)
    _run(para(tf,sa=0,line=1.18),"使うほど、組織の支援が賢くなる。",30,ORANGE,bold=True)
    _run(para(tb(s,0.97,4.7,11.2,1.0),first=True,line=1.4),
     "現場の困りごとを解決した上に、「蓄積したデータが支援を進化させる」将来設計が乗っています。",14,CDW)
hero_slide(s11)

# ============ S12-14 将来設計 ============
image_slide("将来設計 ① 知識循環ループ","diagram_loop.png",
 "毎日の「AIの推論 → 人間の修正 → 成果」を蓄積。これが他社が持ち得ない“データの堀”になる",12,sub="使うほど下書きが現場になじむ")
image_slide("将来設計 ② 記録の蒸留(L1→L4)","diagram_distill.png",
 "断片的・属人的な実践知を、再利用可能な「支援知」へ段階的に蒸留する",13,sub="現場の知恵を資産に")
image_slide("将来設計 ③ 自己改善の二輪(明示 × 暗黙)","diagram_twowheel.png",
 "管理者が定めた基準と、現場の修正から学んだ傾向を、AIの下書きに同時に織り込む",14,sub="施設らしさを生成に反映")

# ============ S15 安心設計 ============
image_slide("安心設計 ― 個人情報を守りながら学ぶ","fig2.png",
 "同意状況の可視化。施設と保護者の同意がそろった記録だけを、個人を特定しない形で活用(個人情報モザイク済)",15,
 sub="同意(施設×児童のAND)・暗号化・k匿名・テナント分離・回数評価をしない設計")

# ============ S16 市場 ============
content_slide("市場と世界の潮流",[
 ("**拡大する市場**: 放課後等デイの費用額は **R6年度 約6,098億円(前年比+14.9%)**、児童発達支援と合わせ **約8,826億円**",0),
 ("利用者 **約30万人**(10年で約5.7倍)、事業所 **全国2万カ所超**。障害福祉全体(約4.18兆円・+12.1%)の中でも**高成長**セグメント",0),
 ("放課後等デイは**文書作成が法令上の必須業務**(計画・モニタリング・記録)。不備は減算・返還リスクに直結",0),
 ("慢性的な人手不足と高い文書負担、記録・支援の属人化。多事業所運営で品質標準化ニーズが拡大",0),
 ("世界的にも**ドキュメント負担の軽減がAIの最有望用途**(医療のAIスクライブ等)。福祉は後続の有望市場",0),
 ("規制の本格化(個人情報保護法・3省2ガイドライン等)= **適合が参入障壁かつ差別化要因**",0),
], 16, size=13.5, gap=7, sub="出典: こども家庭庁/国保連データ(R5→R6年度)・社会福祉施設等調査(R5)")

# ============ S17 強み ============
content_slide("決定的な強み",[
 ("**現場が今すぐ楽になる**: 書類負担を即軽減し、子どもと向き合う時間を取り戻す",0),
 ("**使うほど賢くなるデータの堀**: 修正・理由・成果の蓄積は後追いでは作れない(時間が味方)",0),
 ("**規制適合を内蔵**: 同意・暗号化・匿名化・テナント分離=福祉での“信頼”が堀",0),
 ("**モデル非依存**: OpenAIに縛られず、蓄積で国産・ローカルAIへ移行できる",0),
 ("**育成志向**: 人を置き換えず、観察力・支援設計力を育てる(現場の価値観と一致)",0),
], 17, size=15, gap=11)

# ============ S18 ロードマップ ============
content_slide("ロードマップ",[
 ("**現在(本番稼働)**: 連絡帳・計画書づくりの現場支援、施設記録基準のAI対話、学習基盤(法人内)",0),
 ("**次段**: 介入↔成果の連結 / 外部AI契約(データ保持)整備 / 見本キュレーション運用",0),
 ("**将来**: 匿名化した全国横断の支援知(法務要件クリア後)/ ローカル・国産AIへの移行",0),
 ("導入は無料相談・オンラインデモから。初期設定・データ移行・研修まで伴走",1),
], 18, sub="まず現場価値、その上に将来設計を段階的に")

# ============ S19 クロージング(ヒーロー画像) ============
def s19(s):
    rect(s,0.95,3.45,3.2,0.06,fill=ORANGE)
    tf=tb(s,0.95,1.7,11.5,1.9)
    _run(para(tf,first=True,sa=6,line=1.2),"「書類のための時間」を、",30,WHITE,bold=True)
    _run(para(tf,sa=0,line=1.2),"「子どものための時間」へ。",30,ORANGE,bold=True)
    tf2=tb(s,0.97,3.75,11.4,2.4)
    _run(para(tf2,first=True,sa=8),"きづりは、放課後等デイの記録を「子どもと向き合う時間」と「組織の支援知」に変えます。",15,CDW)
    _run(para(tf2,sa=4),"技術詳細は別添『生成ロジックのアカデミック詳説』をご参照ください。",13,CBW)
    _run(para(tf2,sb=10),"お問い合わせ: soship.co.jp ／ info@soship.co.jp ／ 運営：株式会社ソーシップ",13,CBW)
hero_slide(s19)

out=os.path.join(HERE,"kiduri-紹介資料_商社・共同開発者向け_2026-07-01.pptx")
prs.save(out); print("wrote", out, "slides", len(prs.slides._sldIdLst))
