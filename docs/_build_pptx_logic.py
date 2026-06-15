# -*- coding: utf-8 -*-
"""kiduri 生成ロジック詳説 PowerPoint(技術評価者向け・温かいデザイン)。"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

HERE=os.path.dirname(os.path.abspath(__file__)); FIG=os.path.join(HERE,"figures"); FONT="BIZ UDPGothic"
CREAM=RGBColor(0xF7,0xF5,0xF0); INK=RGBColor(0x2C,0x3E,0x50); ORANGE=RGBColor(0xE8,0x83,0x3A)
TEAL=RGBColor(0x2A,0x8F,0x8F); NAVY=RGBColor(0x1A,0x27,0x44); BEIGE=RGBColor(0xEF,0xEA,0xE0)
GRAY=RGBColor(0x6B,0x7B,0x8D); WHITE=RGBColor(0xFF,0xFF,0xFF); CDW=RGBColor(0xC9,0xD6,0xE6); CBW=RGBColor(0xB8,0xC6,0xD8)
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; SW=13.333; SH=7.5; TOTAL=13

def slide(bg=CREAM):
    s=prs.slides.add_slide(BLANK); r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.shadow.inherit=False; r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); return s
def rect(s,x,y,w,h,fill=None,line=None,shape=MSO_SHAPE.RECTANGLE,radius=None):
    sp=s.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h)); sp.shadow.inherit=False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1.25)
    if radius is not None and shape==MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0]=radius
        except Exception: pass
    return sp
def _run(p,text,size,color=INK,bold=False,italic=False,name=FONT):
    r=p.add_run(); r.text=text; f=r.font; f.size=Pt(size); f.bold=bold; f.italic=italic; f.color.rgb=color; f.name=name
    rPr=r._r.get_or_add_rPr(); ea=rPr.makeelement(qn('a:ea'),{'typeface':name}); rPr.append(ea); return r
def tb(s,x,y,w,h,anchor=MSO_ANCHOR.TOP):
    box=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=box.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor; tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0; return tf
def para(tf,first=False,align=PP_ALIGN.LEFT,sa=8,sb=0,level=0,line=1.18):
    p=tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment=align; p.level=level; p.space_after=Pt(sa); p.space_before=Pt(sb); p.line_spacing=line; return p
def emit(p,text,size,color=INK):
    import re
    for seg in re.split(r'(\*\*[^*]+\*\*)',text):
        if not seg: continue
        if seg.startswith('**'): _run(p,seg[2:-2],size,ORANGE if color==INK else color,bold=True)
        else: _run(p,seg,size,color)
def footer(s,page):
    rect(s,0,7.16,SW,0.018,fill=RGBColor(0xE3,0xDD,0xD2))
    _run(para(tb(s,0.55,7.2,9,0.28),first=True),"きづり 生成ロジック詳説 ｜ 株式会社ソーシップ",9,GRAY)
    _run(para(tb(s,11.2,7.2,1.6,0.28),first=True,align=PP_ALIGN.RIGHT),f"{page} / {TOTAL}",9,GRAY)
def title_bar(s,title,sub=None):
    rect(s,0.55,0.5,0.14,0.66,fill=TEAL,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.5)
    _run(para(tb(s,0.82,0.42,11.9,0.82),first=True,sa=0),title,26,INK,bold=True)
    if sub: _run(para(tb(s,0.84,1.18,11.9,0.4),first=True),sub,13,TEAL)
    rect(s,0.82,1.6 if sub else 1.32,11.7,0.016,fill=RGBColor(0xE3,0xDD,0xD2))
def bullets(s,items,x=0.9,y=1.9,w=11.5,h=5.0,size=14.5,gap=8):
    tf=tb(s,x,y,w,h)
    for i,(text,lvl) in enumerate(items):
        p=para(tf,first=(i==0),level=lvl,sa=gap,line=1.2)
        if lvl==0: _run(p,"▸ ",size+1,ORANGE,bold=True); emit(p,text,size)
        else: _run(p,"・",size-1,TEAL); emit(p,text,size-1,RGBColor(0x4A,0x55,0x60))
    return tf
def formula(s,x,y,w,h,lines):
    sp=rect(s,x,y,w,h,fill=BEIGE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.08)
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    tf.margin_left=Inches(0.3);tf.margin_right=Inches(0.3);tf.margin_top=Inches(0.12);tf.margin_bottom=Inches(0.12)
    for i,t in enumerate(lines):
        p=para(tf,first=(i==0),align=PP_ALIGN.CENTER,sa=4,line=1.25); _run(p,t,15,INK,bold=False)
def image_slide(title,img,caption,page,sub=None):
    s=slide(); title_bar(s,title,sub); footer(s,page)
    path=os.path.join(FIG,img); iw,ih=Image.open(path).size; ar=iw/ih
    bx,by,bw,bh=0.8,1.8,11.73,4.5; w=bw; h=w/ar
    if h>bh: h=bh; w=h*ar
    x=bx+(bw-w)/2; y=by+(bh-h)/2
    rect(s,x-0.05,y-0.05,w+0.1,h+0.1,fill=WHITE,line=RGBColor(0xDD,0xD5,0xC8),shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.02)
    s.shapes.add_picture(path,Inches(x),Inches(y),Inches(w),Inches(h))
    if caption: _run(para(tb(s,0.8,by+bh+0.05,11.73,0.4),first=True,align=PP_ALIGN.CENTER),caption,11,GRAY)
    return s
def content(title,items,page,sub=None,**kw):
    s=slide(); title_bar(s,title,sub); footer(s,page); bullets(s,items,**kw); return s

# S1 タイトル
s=slide(NAVY); rect(s,0,5.0,SW,0.03,fill=TEAL)
_run(para(tb(s,0.95,0.6,9,0.5),first=True),"補足資料(技術評価者向け)",13,CDW)
tf=tb(s,0.9,1.7,11.6,2.4)
_run(para(tf,first=True,sa=2),"生成ロジックの詳説",46,WHITE,bold=True)
_run(para(tf,sa=0,line=1.25),"支援記録・個別支援計画の生成を、形式的に記述する",20,ORANGE,bold=True)
tf2=tb(s,0.92,5.25,11.5,1.4)
_run(para(tf2,first=True,sa=4),"本番コードベースに準拠(準拠コミット d9912d8)",14,WHITE)
_run(para(tf2,sa=0),"2026年7月1日 ／ 運営：株式会社ソーシップ",12.5,CBW)

# S2 全体像
s=slide(); title_bar(s,"全体像 ― 支援文書生成は「条件付き生成」である","児童ごとに集めた文脈 x を条件に、下書き y を生成。確定は必ず人"); footer(s,2)
formula(s,1.4,2.0,10.5,1.0,["y 〜 M_θ( · | Π_s(x) ‖ G(c,d) ; τ )"])
bullets(s,[
 ("**M_θ**=生成モデル(OpenAI gpt-5.4系) / **Π_s**=仮名化 / **G**=施設ガイダンス / **τ**=温度",0),
 ("**不変条件①** 外部モデルへ渡る全テキストは Π_s を通過(静的検査 A005 で強制)",0),
 ("**不変条件②** 各生成は根拠の要約(sources)を返す(説明可能性)",0),
 ("**不変条件③** 生成・修正を学習基盤へ追記(本処理は止めない / 同意は fail-closed)",0),
], y=3.3, h=3.4, size=15, gap=10)

# S3 文脈の階層的集約 + コア目標ロジック
content("個別支援計画の文脈集約と目標ロジック",[
 ("業務データを**決定的に集約**して文脈 x を作る(確率的検索ではない)",0),
 ("土台: **職員アセスメント**＋**保護者アセスメント**(5領域・目標)＋**最新モニタリング**(達成度)",1),
 ("補助: 前回の確定計画 / 能力評価の別添(§8 客観スコア)",1),
 ("**コアの目標ロジック**: 継続目標は**モニタリング**から、新規目標は**職員+保護者アセスメント**から抽出",0),
 ("連絡帳の徹底解析は上流で実施: **アセスメント生成が直近6か月の連絡帳を全件(蒸留)分析**(計画は連絡帳を直接読まない)",0),
 ("どの段を使ったかを sources として返却(説明可能性)",0),
], 3, sub="アセスメント+モニタリングが土台。連絡帳の解析は上流(アセスメント)で吸収される")

# S4 多段生成①
s=slide(); title_bar(s,"多段条件付き生成 ① 個別支援計画・改訂","単一呼出の構造化生成と、低温度の保守的編集"); footer(s,4)
bullets(s,[
 ("**個別支援計画**: gpt-5.4 / τ=0.8 / max 4000 を1回。JSON(意向・方針・長期/短期・5領域明細)を生成",0),
 ("目標は **継続=モニタリング / 新規=職員+保護者アセスメント** を最優先アンカーに(プロンプトで【最重要】指定)",0),
 ("出力健全性: JSONでなければ 502 で明示拒否(壊れた構造は採用しない)→ Π_s⁻¹ で実名復元",0),
 ("**計画改訂**: τ=0.3 の保守的編集。指摘の無い箇所は一字一句保持、変更点は注釈(annotations)で保存",0),
], y=2.0, h=3.0, size=15, gap=10)
formula(s,1.4,5.3,10.5,1.2,["計画改訂: (y₁, A) = M_θ( · | y₀(原案), r(保護者コメント), m(議事録) )","A = 変更注釈 {(field, type∈{追加,削除}, text, reason)}"])

# S5 多段生成②
s=slide(); title_bar(s,"多段条件付き生成 ② アセスメント・連絡帳","逐次の条件付き生成チェーンと、軽量モデル+フォールバック"); footer(s,5)
formula(s,1.0,1.95,11.3,1.55,[
 "u_dom  = M_θ( · | 連絡帳6か月(全件・蒸留), 家庭視点 ; τ=0.6 )   … 5領域課題",
 "u_short = M_θ( · | u_dom, 面談, 家庭視点 ; 0.6 )         … 短期目標",
 "u_long  = M_θ( · | u_short, 面談, 前回長期 ; 0.6 )        … 長期目標"])
bullets(s,[
 ("アセスメントは**3段の連鎖**。各段は前段の出力を条件に取る(現状把握→近接→遠隔目標)",0),
 ("仮名化は連鎖全体で保持。中間生成物は仮名のまま、最終出力でのみ実名復元(外部は実名に触れない)",0),
 ("連絡帳は軽量モデル(gpt-5.4-mini)+ ヒヤリハット検出の副呼出 + 失敗時は決定的連結のフォールバック",0),
], y=3.7, h=3.0, size=14.5, gap=10)

# S6 ガイダンス(図)
image_slide("ガイダンス作用素 ― 明示 ⊕ 暗黙","diagram_twowheel.png",
 "G(c,d) = G_explicit(c) ⊕ [集計同意] · G_implicit(c,d)。明示はPIIなしで常時、暗黙は同意下で作用",6,
 sub="施設の基準と現場の傾向を、下書きへ同時に注入")

# S7 自己改善ループ(図)
s=slide(); title_bar(s,"自己改善ループの形式化","人間の修正を保存し、次期ガイダンスへ"); footer(s,7)
img="diagram_loop.png"; path=os.path.join(FIG,img); iw,ih=Image.open(path).size; ar=iw/ih
w=8.4; h=w/ar; x=(SW-w)/2
rect(s,x-0.05,1.78,w+0.1,h+0.1,fill=WHITE,line=RGBColor(0xDD,0xD5,0xC8),shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.02)
s.shapes.add_picture(path,Inches(x),Inches(1.83),Inches(w),Inches(h))
formula(s,1.0,4.55,11.3,1.6,[
 "δ_k = 1 − similar_text(y_k, y'_k)/100 ∈ [0,1]      ai_acceptance = 1 − mean(δ)",
 "G^(t+1) = Γ( G^(t), {(y, y', 理由, 成果)} )    ※ Γ は M_θ に非依存(=モデル交換可)"])
_run(para(tb(s,1.0,6.3,11.3,0.6),first=True,line=1.25),
 "※ 逆インセンティブ回避: 修正量 δ を職員評価に直結させない(最適化指標と評価指標を分離)",12,GRAY)

# S8 蒸留(図)
image_slide("構造化と蒸留(L1 → L4)","diagram_distill.png",
 "L1 生データ → L2 構造化(本文非保存=PIIなし) → L3 支援知(k匿名 K=5) → L4 原理",8,
 sub="記録を再利用可能な「支援知」へ。支援者成長モデル(Lv1-4)で介入量を適応")

# S9 プライバシーの形式モデル
content("プライバシーの形式モデル",[
 ("**仮名化 Π_s** は既知氏名に対する近似対合(Π_s⁻¹∘Π_s ≈ id)。外部モデルは実名を観測しない",0),
 ("**構造化スクラブ**(日付・電話・番号・敬称付き氏名を決定的除去)+ **短名 fail-safe**(1文字氏名が残る例は破棄)",0),
 ("**同意 AND・fail-closed**: Learn(s) = 施設の集計同意 ∧ 児童の学習同意。未整備の grant は中断",0),
 ("**3層PII**: 原本=暗号化 / 仮名化payload / 非暗号化列に実名を入れない",0),
 ("**テナント分離**(検索・監査を法人スコープに限定)/ **同意文面スナップショット**(立証可能性)",0),
], 9, size=14.5, gap=11)

# S10 成果指標 / 別添「評価状況の全体像」
s=slide(); title_bar(s,"成果(Outcome)― 別添「評価状況の全体像」","個人内評価(他児比較でなく過去の自分との差)。客観 × 本人の主観を統合"); footer(s,10)
bullets(s,[
 ("**個人内評価**: 能力スコアは 0〜10 の個人内基準。他児との比較をしない(福祉の原則)",0),
 ("**A 客観スコア変化**: ability_scores の項目最新値で集計(平均Δ・向上/低下)",0),
 ("**B モニタリング達成度**: 達成度 a∈{1..5} を pct = (mean(a)−1)/4 × 100 で正規化",0),
 ("**C 主観×客観の一致**: 本人の主観(mynameis, 1..5→0..10 正規化)と客観を領域別に突合",0),
 ("別添の客観要約は **個別支援計画AIのプロンプトへ還流**(記録 → 評価 → 計画の循環)",0),
], y=1.95, h=3.7, size=14, gap=8)
formula(s,1.6,5.75,10.1,0.95,["agree_j = max(0, 1 − |o_j − b_j| / 10) × 100      overall = mean_j(agree_j)"])

# S11 説明可能性
content("説明可能性とトレーサビリティ",[
 ("各生成は **sources**(アセス/モニタリング/連絡帳件数と期間/前計画/能力評価の有無)を返す",0),
 ("**AiGenerationLog** に モデル・温度・最大トークン・(マスク済)プロンプト・トークン量・所要時間 を記録",0),
 ("→ 生成の**再現性・監査可能性**を担保。学習基盤は生成(L2)と修正(L1)を追記",0),
 ("これらが将来の Γ(自己改善)とローカルAI学習の素材になる",0),
], 11, size=15, gap=12, sub="「何を根拠に、どのモデルで、どう生成したか」を後から検証できる")

# S12 限界と拡張
content("限界と理論的拡張(誠実な開示)",[
 ("**RAG は生成へ未配線**: 埋め込み基盤は実装済だが、生成は決定的に集約した文脈に条件付け(検索挿入は将来)",0),
 ("**Γ は現状ヒューリスティック**(統計写像)。蓄積は将来のローカルAI蒸留の素材として保存",0),
 ("**因果の留保**: 成果は相関・記述統計。支援→成果の因果は S9(介入↔成果連結)で前提整備",0),
 ("**全国横断は未解禁**: 法務4点+恒久匿名化のクリアを前提にゲート",0),
 ("**モデル非依存の含意**: Γ と蓄積が M_θ と独立 = 国産・ローカルAIへ置換しても機能(事業命題の技術的裏づけ)",0),
], 12, size=14.5, gap=10)

# S13 クロージング
s=slide(NAVY); rect(s,0.95,3.3,3.2,0.06,fill=TEAL)
tf=tb(s,0.95,1.9,11.5,1.6)
_run(para(tf,first=True,sa=4,line=1.2),"AIは交換可能、",30,WHITE,bold=True)
_run(para(tf,sa=0,line=1.2),"蓄積した支援知は交換できない。",30,ORANGE,bold=True)
tf2=tb(s,0.97,3.6,11.4,2.2)
_run(para(tf2,first=True,sa=8),"生成は『仮名化された条件付き生成』、改善は『人間の修正を保存して学ぶループ』。",14.5,CDW)
_run(para(tf2,sa=4),"いずれもモデルに依存しない設計で、将来のローカルAIへ継承できます。",14.5,CDW)
_run(para(tf2,sb=10),"概要は別冊『紹介資料』をご参照ください。 ／ 運営：株式会社ソーシップ",13,CBW)

out=os.path.join(HERE,"kiduri-補足_生成ロジック詳説_2026-07-01.pptx")
prs.save(out); print("wrote",out,"slides",len(prs.slides._sldIdLst))
